#!/usr/bin/env python3
"""
PPTX → HTML スライドコンバーター
各PPTXファイルを個別HTMLページに変換する
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import html
import glob
import json

REPO_DIR = os.path.dirname(os.path.abspath(__file__))
PPTX_DIR = "/mnt/c/Users/konog/.claude/URIKAIコンテンツ"
SLIDES_DIR = os.path.join(REPO_DIR, "slides")
os.makedirs(SLIDES_DIR, exist_ok=True)


def rgb_to_css(rgb_color):
    """RGBColor to CSS hex"""
    if rgb_color is None:
        return None
    try:
        return f"#{rgb_color.red:02x}{rgb_color.green:02x}{rgb_color.blue:02x}"
    except:
        return None


def get_fill_color(shape):
    """Get shape fill color"""
    try:
        fill = shape.fill
        if fill.type is not None:
            try:
                return rgb_to_css(fill.fore_color.rgb)
            except:
                pass
    except:
        pass
    return None


def get_slide_bg_color(slide):
    """Get slide background color"""
    try:
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            try:
                return rgb_to_css(fill.fore_color.rgb)
            except:
                pass
    except:
        pass
    return "#f4f6f7"


def process_text_frame(tf):
    """Convert text frame to HTML"""
    parts = []
    for para in tf.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Get paragraph styling
        font_size = None
        font_color = None
        is_bold = False
        align = "left"

        if para.font:
            if para.font.size:
                font_size = para.font.size.pt
            font_color = rgb_to_css(para.font.color.rgb) if para.font.color and para.font.color.rgb else None
            is_bold = para.font.bold or False

        if para.alignment == PP_ALIGN.CENTER:
            align = "center"
        elif para.alignment == PP_ALIGN.RIGHT:
            align = "right"

        # Build style
        styles = []
        if font_size:
            # Scale for web (larger for readability)
            if font_size >= 26:
                # Titles - keep current size
                web_size = font_size * 0.85
            else:
                # Body text - make significantly larger
                web_size = max(font_size * 1.15, 16)
            styles.append(f"font-size:{web_size:.0f}px")
        if font_color:
            styles.append(f"color:{font_color}")
        if is_bold:
            styles.append("font-weight:700")
        if align != "left":
            styles.append(f"text-align:{align}")

        style_str = ";".join(styles)
        escaped = html.escape(text)
        parts.append(f'<p style="{style_str}">{escaped}</p>')

    return "\n".join(parts)


def shape_to_html(shape):
    """Convert a shape to HTML"""
    result = ""

    # Get position and size
    left = shape.left / 914400 if shape.left else 0  # EMU to inches
    top = shape.top / 914400 if shape.top else 0
    width = shape.width / 914400 if shape.width else 0
    height = shape.height / 914400 if shape.height else 0

    # Scale to percentage of slide (13.333 x 7.5 inches)
    left_pct = (left / 13.333) * 100
    top_pct = (top / 7.5) * 100
    width_pct = (width / 13.333) * 100
    height_pct = (height / 7.5) * 100

    fill_color = get_fill_color(shape)

    # Build shape div
    styles = [
        f"position:absolute",
        f"left:{left_pct:.2f}%",
        f"top:{top_pct:.2f}%",
        f"width:{width_pct:.2f}%",
        f"height:{height_pct:.2f}%",
        "overflow:hidden",
    ]

    if fill_color:
        styles.append(f"background:{fill_color}")

    # Check for rounded rectangle
    try:
        from pptx.enum.shapes import MSO_SHAPE
        if hasattr(shape, 'shape_type') and shape.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
            styles.append("border-radius:8px")
        elif hasattr(shape, 'shape_type') and shape.auto_shape_type == MSO_SHAPE.OVAL:
            styles.append("border-radius:50%")
    except:
        pass

    style_str = ";".join(styles)

    # Process text
    text_html = ""
    if shape.has_text_frame:
        text_html = process_text_frame(shape.text_frame)

    # Process table
    if shape.has_table:
        text_html = table_to_html(shape.table)

    if text_html or fill_color:
        result = f'<div style="{style_str}">{text_html}</div>'

    return result


def table_to_html(table):
    """Convert table to HTML"""
    rows_html = []
    for row_idx, row in enumerate(table.rows):
        cells_html = []
        for cell in row.cells:
            text = cell.text.strip()
            escaped = html.escape(text)

            if row_idx == 0:
                cells_html.append(
                    f'<th style="background:#1a3c6e;color:#fff;padding:8px 12px;'
                    f'font-size:18px;text-align:left;border:1px solid #30363d">{escaped}</th>'
                )
            else:
                bg = "#eaecee" if row_idx % 2 == 0 else "transparent"
                cells_html.append(
                    f'<td style="background:{bg};padding:8px 12px;font-size:18px;'
                    f'color:#2c3e50;border:1px solid #30363d">{escaped}</td>'
                )
        tag = "tr"
        rows_html.append(f"<{tag}>{''.join(cells_html)}</{tag}>")

    return (
        '<table style="border-collapse:collapse;width:100%;margin-top:4px">'
        + "".join(rows_html)
        + "</table>"
    )


def slide_to_html(slide, slide_num):
    """Convert one slide to HTML div"""
    bg_color = get_slide_bg_color(slide)

    shapes_html = []
    for shape in slide.shapes:
        h = shape_to_html(shape)
        if h:
            shapes_html.append(h)

    return (
        f'<div class="slide" data-slide="{slide_num}" style="background:{bg_color}">'
        + "\n".join(shapes_html)
        + "</div>"
    )


def convert_pptx_to_html(pptx_path, output_path, section_num, title):
    """Convert a PPTX file to a standalone HTML viewer page"""
    prs = Presentation(pptx_path)
    total_slides = len(prs.slides)

    slides_html = []
    for i, slide in enumerate(prs.slides):
        slides_html.append(slide_to_html(slide, i + 1))

    all_slides = "\n".join(slides_html)

    page_html = f'''<!DOCTYPE html>
<html lang="ja">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>URIKAI - {html.escape(title)}</title>
<style>
:root {{
    --bg: #0d1117;
    --bg-card: #161b22;
    --accent: #d4a537;
    --text: #e6edf3;
    --text-muted: #8b949e;
    --border: #30363d;
}}
* {{ margin:0; padding:0; box-sizing:border-box; }}
body {{
    font-family: 'Segoe UI','Meiryo UI',sans-serif;
    background: var(--bg);
    color: var(--text);
    height: 100vh;
    display: flex;
    flex-direction: column;
    user-select: none;
    -webkit-user-select: none;
}}
.toolbar {{
    background: var(--bg-card);
    border-bottom: 2px solid var(--accent);
    padding: 0.5rem 1.5rem;
    display: flex;
    align-items: center;
    gap: 1rem;
    flex-shrink: 0;
}}
.toolbar a {{
    color: var(--accent);
    text-decoration: none;
    font-size: 0.9rem;
    font-weight: 600;
}}
.toolbar a:hover {{ text-decoration: underline; }}
.toolbar .title {{
    color: var(--text);
    font-size: 0.95rem;
    font-weight: 600;
    flex: 1;
}}
.nav-controls {{
    display: flex;
    align-items: center;
    gap: 0.5rem;
}}
.nav-btn {{
    background: var(--border);
    color: var(--text);
    border: none;
    width: 36px;
    height: 36px;
    border-radius: 6px;
    font-size: 1.1rem;
    cursor: pointer;
    display: flex;
    align-items: center;
    justify-content: center;
    transition: background 0.15s;
}}
.nav-btn:hover {{ background: var(--accent); color: #0a2f5c; }}
.nav-btn:disabled {{ opacity: 0.3; cursor: default; }}
.nav-btn:disabled:hover {{ background: var(--border); color: var(--text); }}
.page-info {{
    color: var(--text-muted);
    font-size: 0.8rem;
    min-width: 60px;
    text-align: center;
}}
.slide-container {{
    flex: 1;
    display: flex;
    align-items: center;
    justify-content: center;
    padding: 1rem;
    overflow: hidden;
}}
.slide {{
    width: 100%;
    max-width: 1400px;
    aspect-ratio: 16/9;
    position: relative;
    border-radius: 6px;
    box-shadow: 0 4px 30px rgba(0,0,0,0.5);
    display: none;
    overflow: hidden;
}}
.slide.active {{
    display: block;
}}
.slide p {{
    margin: 0;
    padding: 1px 0;
    line-height: 1.5;
}}
.slide table {{
    font-family: 'Segoe UI','Meiryo UI',sans-serif;
}}
/* Keyboard hint */
.hint {{
    text-align: center;
    padding: 0.4rem;
    color: var(--text-muted);
    font-size: 0.7rem;
    flex-shrink: 0;
}}
</style>
</head>
<body>
<div class="toolbar">
    <a href="../index.html">&larr; Back</a>
    <span class="title">Section {section_num} — {html.escape(title)}</span>
    <div class="nav-controls">
        <button class="nav-btn" id="prevBtn" onclick="navigate(-1)">&larr;</button>
        <span class="page-info" id="pageInfo">1 / {total_slides}</span>
        <button class="nav-btn" id="nextBtn" onclick="navigate(1)">&rarr;</button>
    </div>
</div>
<div class="slide-container">
{all_slides}
</div>
<div class="hint">&larr; &rarr; Arrow keys to navigate</div>
<script>
let current = 1;
const total = {total_slides};
function navigate(dir) {{
    current = Math.max(1, Math.min(total, current + dir));
    document.querySelectorAll('.slide').forEach(s => s.classList.remove('active'));
    document.querySelector('[data-slide="'+current+'"]').classList.add('active');
    document.getElementById('pageInfo').textContent = current + ' / ' + total;
    document.getElementById('prevBtn').disabled = current === 1;
    document.getElementById('nextBtn').disabled = current === total;
}}
// Init
document.querySelector('[data-slide="1"]').classList.add('active');
document.getElementById('prevBtn').disabled = true;
// Keyboard
document.addEventListener('keydown', e => {{
    if (e.key === 'ArrowRight' || e.key === ' ') navigate(1);
    if (e.key === 'ArrowLeft') navigate(-1);
}});
// Block right-click
document.addEventListener('contextmenu', e => e.preventDefault());
// Auth check
if (sessionStorage.getItem('urikai_auth') !== 'true') {{
    window.location.href = '../index.html';
}}
</script>
</body>
</html>'''

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(page_html)

    return total_slides


def main():
    sections = [
        ("01", "金融市場とは何か"), ("02", "金融市場の種類"), ("03", "市場参加者"),
        ("04", "通貨ペアの基礎"), ("05", "市場の仕組み"), ("06", "注文の種類"),
        ("07", "取引セッションと時間帯"), ("08", "ローソク足パターン"),
        ("09", "サポート&レジスタンス"), ("10", "テクニカルインジケーター"),
        ("11", "チャートパターン"), ("12", "フィボナッチ"),
        ("13", "マルチタイムフレーム分析"), ("14", "エリオット波動"),
        ("15", "ハーモニックパターン"), ("16", "ワイコフ理論"),
        ("17", "オーダーフロー分析"), ("18", "ボリュームプロファイル"),
        ("19", "インターマーケット分析"), ("20", "マクロ経済指標"),
        ("21", "中央銀行政策"), ("22", "地政学リスクとニューストレード"),
        ("23", "各国経済プロファイル"), ("24", "COTレポートとポジション分析"),
        ("25", "センチメント分析"), ("26", "ポジションサイジング"),
        ("27", "リスクリワード比と期待値"), ("28", "SLとTP戦略"),
        ("29", "ドローダウン管理"), ("30", "感情管理"),
        ("31", "認知バイアス"), ("32", "規律とルーティン"),
        ("33", "トレーディングプラン作成"), ("34", "バックテストとシステム評価"),
        ("35", "トレード戦略集"),
    ]

    # Map filenames
    pptx_files = glob.glob(os.path.join(PPTX_DIR, "*.pptx"))
    file_map = {}
    for f in pptx_files:
        basename = os.path.basename(f)
        num = basename.split("_")[0]
        file_map[num] = f

    print("=" * 60)
    print("PPTX → HTML 変換")
    print("=" * 60)

    total_slides = 0
    converted = 0

    for num, title in sections:
        if num not in file_map:
            print(f"  SKIP: {num} {title} (PPTX not found)")
            continue

        pptx_path = file_map[num]
        output_path = os.path.join(SLIDES_DIR, f"{num}.html")

        try:
            slides = convert_pptx_to_html(pptx_path, output_path, num, title)
            total_slides += slides
            converted += 1
            print(f"  OK: {num}_{title}.html ({slides} slides)")
        except Exception as e:
            print(f"  ERROR: {num} {title}: {e}")

    print(f"\n  変換完了: {converted}ファイル / {total_slides}スライド")
    print(f"  出力先: {SLIDES_DIR}/")


if __name__ == "__main__":
    main()
